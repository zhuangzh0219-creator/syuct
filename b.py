
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import math


prs = Presentation("C:/Users/zhuan_rzwfs19/tmp/work/work08/aaa/a.pptx")


# =========================================
# 工具函数
# =========================================

def center(shape):
    """
    shape中心点
    """
    return (
        shape.left + shape.width / 2,
        shape.top + shape.height / 2
    )


def distance(x1, y1, x2, y2):
    return math.sqrt(
        (x1 - x2) ** 2 +
        (y1 - y2) ** 2
    )


def nearest_node(x, y, nodes):
    """
    找离某个点最近的节点
    """

    nearest = None
    min_dist = 999999999

    for node in nodes:

        cx, cy = center(node["shape"])

        d = distance(x, y, cx, cy)

        if d < min_dist:
            min_dist = d
            nearest = node

    return nearest


# =========================================
# 主逻辑
# =========================================

slides_result = []

for slide_idx, slide in enumerate(prs.slides, start=1):

    print("\n====================")
    print("slide:", slide_idx)

    nodes = []
    edges = []

    # =====================================
    # 第一轮：收集节点
    # =====================================

    for shape in slide.shapes:

        # -----------------------------
        # 文本/图形节点
        # -----------------------------
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:

            node = {
                "id": f"node_{shape.shape_id}",
                "shape": shape,
                "text": "",
                "shape_type": "",
                "x": shape.left,
                "y": shape.top,
                "w": shape.width,
                "h": shape.height
            }

            # 图形具体类型
            try:
                node["shape_type"] = str(
                    shape.auto_shape_type
                )
            except:
                node["shape_type"] = "UNKNOWN"

            # 文本
            if hasattr(shape, "text"):
                node["text"] = shape.text.strip()

            nodes.append(node)

    # =====================================
    # 第二轮：收集线
    # =====================================

    for shape in slide.shapes:

        if shape.shape_type == MSO_SHAPE_TYPE.LINE:

            print("\n------ LINE FOUND ------")

            print("shape id:", shape.shape_id)

            # =================================
            # 起点终点
            # =================================

            start_x = shape.begin_x
            start_y = shape.begin_y

            end_x = shape.end_x
            end_y = shape.end_y

            print("start:", start_x, start_y)
            print("end:", end_x, end_y)

            # =================================
            # 箭头类型
            # =================================

            xml = shape.element.xml

            arrow_start = False
            arrow_end = False

            if "a:headEnd" in xml:
                arrow_end = True

            if "a:tailEnd" in xml:
                arrow_start = True

            print("arrow_start:", arrow_start)
            print("arrow_end:", arrow_end)

            # =================================
            # 匹配最近节点
            # =================================

            from_node = nearest_node(
                start_x,
                start_y,
                nodes
            )

            to_node = nearest_node(
                end_x,
                end_y,
                nodes
            )

            edge = {
                "from": from_node["id"]
                    if from_node else None,

                "to": to_node["id"]
                    if to_node else None,

                "from_text": from_node["text"]
                    if from_node else None,

                "to_text": to_node["text"]
                    if to_node else None,

                "arrow_start": arrow_start,
                "arrow_end": arrow_end
            }

            edges.append(edge)

    # =====================================
    # 输出
    # =====================================

    print("\n===== NODES =====")
    pprint(nodes)

    print("\n===== EDGES =====")
    pprint(edges)

    slides_result.append({
        "slide": slide_idx,
        "nodes": nodes,
        "edges": edges
    })


# =========================================
# Mermaid 生成示例
# =========================================

print("\n\n======================")
print("MERMAID")
print("======================")

for slide_data in slides_result:

    print("\nflowchart LR")

    # -----------------------------
    # node
    # -----------------------------

    for node in slide_data["nodes"]:

        node_id = node["id"]
        text = node["text"]

        shape_type = node["shape_type"]

        # 数据库
        if "MAGNETIC_DISK" in shape_type:

            print(
                f'{node_id}[("{text}")]'
            )

        # 普通流程
        elif "PROCESS" in shape_type:

            print(
                f'{node_id}["{text}"]'
            )

        # 普通矩形
        else:

            print(
                f'{node_id}["{text}"]'
            )

    # -----------------------------
    # edge
    # -----------------------------

    for edge in slide_data["edges"]:

        if edge["arrow_end"]:

            print(
                f'{edge["from"]} --> {edge["to"]}'
            )

        else:

            print(
                f'{edge["from"]} --- {edge["to"]}'
            )