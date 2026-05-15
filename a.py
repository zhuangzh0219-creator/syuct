import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
# ppt_path = "C:/Users/tR16277/Desktop/work/work06/file/⓪【リボ利率見直し】PJ計画書_v1.0.pptx"
# ppt_path = "C:/Users/tR16277/Desktop/work/work06/file/② 基本設計書（対応方針）_SR・CR新利率体系_ver1.2.0.pptx"
ppt_path = "C:/Users/zhuan_rzwfs19/tmp/work/work08/aaa/a.pptx"
prs = pptx.Presentation(ppt_path)


def get_strike(run):
    rPr = run._r.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
    )

    if rPr is None:
        return False

    strike = rPr.get("strike")

    return strike not in (None, "noStrike")


def is_hidden_merge_cell(cell):
    tc = cell._tc

    return (
        tc.get("hMerge") == "1"
        or tc.get("vMerge") == "1"
    )


def point_in_rect(px, py, left, top, width, height):

    return (
        left <= px <= left + width
        and top <= py <= top + height
    )


def iter_table_cells(shape):

    table = shape.table

    current_y = shape.top

    for r, row in enumerate(table.rows):

        row_height = row.height
        current_x = shape.left

        for c, cell in enumerate(row.cells):

            col_width = table.columns[c].width

            yield {
                "row": r,
                "col": c,
                "cell": cell,
                "left": current_x,
                "top": current_y,
                "width": col_width,
                "height": row_height,
            }

            current_x += col_width

        current_y += row_height


def extract_cell_text(cell):

    text_parts = []

    for para in cell.text_frame.paragraphs:

        para_text = ""

        for run in para.runs:

            strike = get_strike(run)

            if strike:
                para_text += f"<del:{run.text}>"
            else:
                para_text += run.text

        text_parts.append(para_text)

    return "\n".join(text_parts).strip()


if __name__ == "__main__":

    for slide_idx, slide in enumerate(prs.slides):

        print("=" * 80)
        print(f"SLIDE {slide_idx + 1}")
        print("=" * 80)

        # ==========================================================
        # 1. 输出所有表格
        # ==========================================================

        for shape_idx, shape in enumerate(slide.shapes):

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:

                print(f"\n[TABLE {shape_idx}]")

                table = shape.table

                for r, row in enumerate(table.rows):

                    row_data = []

                    for c, cell in enumerate(row.cells):

                        if is_hidden_merge_cell(cell):
                            continue

                        cell_text = extract_cell_text(cell)

                        row_data.append(repr(cell_text))

                    print(row_data)

        # ==========================================================
        # 2. 查找 CALLOUT
        # ==========================================================

        for shape_idx, shape in enumerate(slide.shapes):

            try:
                auto_type = str(shape.auto_shape_type)
            except:
                continue

            auto_type = str(shape.auto_shape_type)

            if "CALLOUT" not in auto_type:
                continue

            callout_text = shape.text.strip()

            print("\n")
            print("=" * 40)
            print("发现 CALLOUT")
            print("=" * 40)

            print("文本:", repr(callout_text))
            print("类型:", auto_type)

            tail_x = shape.left
            tail_y = shape.top + shape.height

            print("tail_x:", tail_x)
            print("tail_y:", tail_y)

            # ======================================================
            # 3. 查找尾巴指向哪个 cell
            # ======================================================

            found = False

            for s in slide.shapes:

                if s.shape_type != MSO_SHAPE_TYPE.TABLE:
                    continue

                for info in iter_table_cells(s):

                    if point_in_rect(
                        tail_x,
                        tail_y,
                        info["left"],
                        info["top"],
                        info["width"],
                        info["height"]
                    ):

                        target_text = extract_cell_text(info["cell"])

                        print("\n>>> 指向 Table Cell <<<")
                        print("row:", info["row"])
                        print("col:", info["col"])
                        print("text:", repr(target_text))

                        found = True

            if not found:
                print("\n!!! 没找到对应 cell !!!")