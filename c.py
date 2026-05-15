from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

items = [
    {
        'type': 'text',
        'text': '３．開発概要　（１）新規入会時の利率判定ランクに 応じた利率設定＜SR＞\n',
        'left': 537281,
        'top': 128501,
        'width': 8988226,
        'height': 394891
    },
    {
        'type': MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        'text': '\n　　スマート入会で設定する不備パラメータをもとに、サミットホストで契約利率を設定。\n以下①～⑤の順で、③を2案で検討。(詳細は５．システム対応案 評価参照)\n　　　\n　　　②スマート入会で不備理由コードを設定し、ホストへ連携。 （開発不要）\n　　　③本登録で不備理由コードに紐づく利率判定ランクを設定した後、利率判定ランクに紐づく利率を設定。\n　　　④送付台紙、計算書等、各種書面に②で求めた利率を表示。\n　　　\nスマート入会とホストで持つ不備パラメータ⇔利率の紐づけは平仄をとる必要がある。\n',
        'left': 165100,
        'top': 867256,
        'width': 9575800,
        'height': 2536899
    },
    {
        'type': MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        'text': '\n',
        'left': 165100,
        'top': 3771204,
        'width': 9575800,
        'height': 2967534
    },
    {
        'type': 'text',
        'text': '■検討経緯、開発概要\n',
        'left': 282101,
        'top': 568363,
        'width': 3521414,
        'height': 338554
    }
]

for item in items:

    if item["type"] == "text":

        shape = slide.shapes.add_textbox(
            Emu(item["left"]),
            Emu(item["top"]),
            Emu(item["width"]),
            Emu(item["height"]),
        )

    else:

        shape = slide.shapes.add_shape(
            item["type"],
            Emu(item["left"]),
            Emu(item["top"]),
            Emu(item["width"]),
            Emu(item["height"]),
        )

    tf = shape.text_frame
    tf.clear()

    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = item["text"]

    p.space_before = 0
    p.space_after = 0
    p.line_spacing = 1.0
    p.alignment = PP_ALIGN.LEFT

    for run in p.runs:

        font = run.font

        # 日本企业PPT关键字体
        font.name = "MS PGothic"

        # 主标题
        if item["top"] < 300000:
            font.size = Pt(18)
            font.bold = True

        # 小标题
        elif item["top"] < 800000:
            font.size = Pt(12)
            font.bold = True

        # 正文
        else:
            font.size = Pt(9)

prs.save("output.pptx")
print("done")