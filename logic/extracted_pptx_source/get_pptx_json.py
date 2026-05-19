#!/usr/bin/env python3
"""
PPTX → Structured JSON Parser
解析 PowerPoint 文件中的所有元素，输出可供 Markdown/LLM 读取的结构化 JSON。

重点解析：
 - 图形（Shape）及其内部文本（保留删除线）
 - 图形间的逻辑指向线（Connector / Arrow）
 - 表格（Table），含被图形覆盖时的层叠关系
 - 图片（Picture）
 - 分组（Group）递归展开
 - z-order 层叠关系（用于判断覆盖逻辑）
"""

import sys
import json
from pathlib import Path
from typing import Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


# ─────────────────────────── 工具函数 ─────────────────────────── 
def emu_to_pt(emu: int | None) -> float | None:
    """EMU → points（保留 2 位小数）"""
    return round(emu / 12700, 2) if emu is not None else None


def bbox(shape) -> dict:
    """返回形状边界框（单位 pt）"""
    return {
        "left":   emu_to_pt(shape.left),
        "top":    emu_to_pt(shape.top),
        "width":  emu_to_pt(shape.width),
        "height": emu_to_pt(shape.height),
    }


def bbox_overlaps(a: dict, b: dict) -> bool:
    """判断两个 bbox 是否重叠"""
    if None in (a["left"], a["top"], b["left"], b["top"]):
        return False
    ax2 = a["left"] + (a["width"] or 0)
    ay2 = a["top"]  + (a["height"] or 0)
    bx2 = b["left"] + (b["width"] or 0)
    by2 = b["top"]  + (b["height"] or 0)
    return a["left"] < bx2 and ax2 > b["left"] and a["top"] < by2 and ay2 > b["top"]


# ─────────────────────────── 文本解析 ───────────────────────────
def parse_run(run) -> dict | None:
    """解析一个文字 run，提取文本与删除线信息（忽略颜色）"""
    text = run.text
    if not text:
        return None
    result: dict[str, Any] = {"text": text}

    # ── 删除线检测（直接读 a:rPr 的 strike 属性，最可靠）── 
    strike = False
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        s_val = rPr.get("strike", "noStrike")
        strike = s_val not in ("noStrike", "")
    if strike:
        result["strikethrough"] = True

    # ── 粗体 / 斜体 ──
    font = run.font
    if font.bold:
        result["bold"] = True
    if font.italic:
        result["italic"] = True

    return result


def parse_paragraph(para) -> dict | None:
    """解析段落，返回 runs 列表 + 对齐信息"""
    runs = [r for r in (parse_run(run) for run in para.runs) if r]
    # 合并相邻无特殊属性的 run
    merged: list[dict] = []
    for r in runs:
        if merged and set(merged[-1].keys()) == {"text"} and set(r.keys()) == {"text"}:
            merged[-1]["text"] += r["text"]
        else:
            merged.append(r)
    if not merged:
        return None
    result: dict[str, Any] = {"runs": merged}
    align = para.alignment
    if align is not None:
        result["align"] = str(align).split(".")[-1].lower()
    return result


def parse_text_frame(tf) -> list[dict]:
    """解析文本框，返回段落列表（保留删除线，忽略颜色）"""
    paras = [p for p in (parse_paragraph(para) for para in tf.paragraphs) if p]
    return paras


def text_plain(tf) -> str:
    """返回文本框纯文本（用于摘要）"""
    lines = []
    for para in tf.paragraphs:
        line = "".join(r.text for r in para.runs)
        if line:
            lines.append(line)
    return "\n".join(lines)


# ─────────────────────────── 连接线解析 ─────────────────────────── 
def is_connector(shape) -> bool:
    """判断是否为连接线/箭头"""
    try:
        st = shape.shape_type
        # MSO_SHAPE_TYPE.LINE = 9
        if st == MSO_SHAPE_TYPE.LINE:
            return True
    except Exception:
        pass
    # 通过 XML tag 判断
    tag = shape._element.tag.split("}")[-1] if "}" in shape._element.tag else shape._element.tag
    if tag in ("cxnSp",):
        return True
    # spPr 下有 prstGeom 的线型
    try:
        sp_pr = shape._element.find(".//" + qn("p:spPr"))
        if sp_pr is None:
            sp_pr = shape._element.find(".//" + qn("a:spPr"))
        if sp_pr is not None:
            geom = sp_pr.find(qn("a:prstGeom"))
            if geom is not None:
                prst = geom.get("prst", "")
                if "Arrow" in prst or prst in ("line", "straightConnector1",
                                               "bentConnector2", "bentConnector3",
                                               "curvedConnector2", "curvedConnector3"):
                    return True
    except Exception:
        pass
    return False


def parse_connector(shape, shape_id_map: dict) -> dict:
    """解析连接线，提取起点/终点形状引用"""
    info: dict[str, Any] = {
        "id":   shape.shape_id,
        "name": shape.name,
        "type": "connector",
        "bbox": bbox(shape),
    }
    # 解析 XML 中的 stCxn / endCxn 连接关系
    elem = shape._element
    for cxn_tag, key in [("a:stCxn", "from_shape"), ("a:endCxn", "to_shape")]:
        cxn = elem.find(".//" + qn(cxn_tag))
        if cxn is not None:
            sid = int(cxn.get("id", 0))
            info[key] = {
                "shape_id": sid,
                "shape_name": shape_id_map.get(sid, {}).get("name"),
                "idx": int(cxn.get("idx", 0)),
            }
    # 箭头方向（从 ln 标签读取）
    arrow_info: dict[str, str] = {}
    for end, tag in [("tail", "a:tailEnd"), ("head", "a:headEnd")]:
        ae = elem.find(".//" + qn(tag))
        if ae is not None:
            t = ae.get("type", "none")
            if t != "none":
                arrow_info[end] = t
    if arrow_info:
        info["arrows"] = arrow_info
    # 文本（连接线偶尔有标注文字）
    try:
        if shape.has_text_frame:
            txt = text_plain(shape.text_frame).strip()
            if txt:
                info["label"] = txt
    except Exception:
        pass
    return info


# ─────────────────────────── 表格解析 ───────────────────────────
def is_hidden_merge_cell(cell):
    tc = cell._tc

    return (
        tc.get("hMerge") is not None
        or tc.get("vMerge") is not None
    )

def parse_table(shape) -> dict:
    info: dict[str, Any] = {
        "id":   shape.shape_id,
        "name": shape.name,
        "type": "table",
        "bbox": bbox(shape),
        "rows": [],
    }
    tbl = shape.table
    info["row_count"] = len(list(tbl.rows))
    info["col_count"] = len(list(tbl.columns))

    raw_col_widths = [col.width for col in tbl.columns]
    raw_row_heights = [row.height for row in tbl.rows]
    scale_x = shape.width / sum(raw_col_widths)
    scale_y = shape.height / sum(raw_row_heights)
    col_widths = [w * scale_x for w in raw_col_widths]
    row_heights = [h * scale_y for h in raw_row_heights]
    table_left = shape.left
    table_top = shape.top

    def find_master(r_idx, c_idx, direction):
        """向上(v)或向左(h)找母单元格行/列索引"""
        if direction == "v":
            r = r_idx
            while r > 0:
                r -= 1  # ← 先减，再判断
                tc = tbl.rows[r].cells[c_idx]._tc
                if tc.get("vMerge") is None:
                    return r
            return 0
        else:
            c = c_idx
            while c > 0:
                c -= 1
                tc = tbl.rows[r_idx].cells[c]._tc
                if tc.get("hMerge") is None:
                    return c
            return 0

    def parse_cell(cell, r_idx: int, c_idx: int) -> dict:
        cell_info: dict[str, Any] = {"row": r_idx, "col": c_idx}
        tc = cell._tc

        is_v_placeholder = tc.get("vMerge") is not None
        is_h_placeholder = tc.get("hMerge") is not None

        if is_v_placeholder or is_h_placeholder:
            # 影子セル: master_cell を記録して内容は空
            master_r = find_master(r_idx, c_idx, "v") if is_v_placeholder else r_idx
            master_c = find_master(r_idx, c_idx, "h") if is_h_placeholder else c_idx
            return {
                "row": r_idx, "col": c_idx,
                "col_span": 1, "row_span": 1,
                "master_cell": {"row": master_r, "col": master_c},
                "content": []
            }

        # 母セルまたは独立セル: XML から span を直接読む
        grid_span = int(tc.get("gridSpan", 1))
        row_span  = int(tc.get("rowSpan",  1))

        cell_info["col_span"] = grid_span
        cell_info["row_span"] = row_span

        left   = table_left + sum(col_widths[:c_idx])
        top    = table_top  + sum(row_heights[:r_idx])
        width  = sum(col_widths[c_idx  : c_idx  + grid_span])
        height = sum(row_heights[r_idx : r_idx  + row_span])

        cell_info["bbox"] = {
            "left":   emu_to_pt(left),
            "top":    emu_to_pt(top),
            "width":  emu_to_pt(width),
            "height": emu_to_pt(height),
        }

        try:
            paras = parse_text_frame(cell.text_frame)
            cell_info["content"] = paras if paras else []
        except Exception:
            cell_info["content"] = []

        return cell_info

    for r_idx, row in enumerate(tbl.rows):
        row_data = []
        for c_idx, cell in enumerate(row.cells):
            if is_hidden_merge_cell(cell):
                # 跳过被合并隐藏的单元格
                continue
            row_data.append(parse_cell(cell, r_idx, c_idx))
        info["rows"].append(row_data)
    return info


# ─────────────────────────── 图片解析 ───────────────────────────

def parse_picture(shape) -> dict:
    """解析图片形状"""
    info: dict[str, Any] = {
        "id":   shape.shape_id,
        "name": shape.name,
        "type": "picture",
        "bbox": bbox(shape),
    }
    try:
        img = shape.image
        info["image_format"] = img.content_type
        # 文件名/blob hash
        info["image_blob_hash"] = hash(img.blob) & 0xFFFFFFFF
    except Exception:
        pass
    # 替换文字（alt text）
    try:
        desc = shape._element.find(".//" + qn("p:nvPicPr"))
        if desc is not None:
            nvpr = desc.find(qn("p:nvPr"))
            if nvpr is not None:
                ph = nvpr.find(qn("p:ph"))
                if ph is None:
                    # 尝试 cNvPr descr
                    cnvpr = desc.find(qn("p:cNvPr"))
                    if cnvpr is not None:
                        descr = cnvpr.get("descr", "")
                        if descr:
                            info["alt_text"] = descr
    except Exception:
        pass
    return info


# ─────────────────────────── 普通形状解析 ───────────────────────────

def get_shape_geometry(shape) -> str | None:
    """获取预设几何形状名称"""
    try:
        sp_pr = shape._element.find(".//" + qn("p:spPr"))
        if sp_pr is None:
            sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is not None:
            geom = sp_pr.find(qn("a:prstGeom"))
            if geom is not None:
                return geom.get("prst")
    except Exception:
        pass
    return None


def get_shape_adjustments(shape) -> dict:
    """
    读取 shape 的 adjustment 参数（尤其用于吹き出し）
    """
    result = {}

    try:
        geom = shape._element.find(".//" + qn("a:prstGeom"))
        if geom is None:
            return result

        avLst = geom.find(qn("a:avLst"))
        if avLst is None:
            return result

        for gd in avLst.findall(qn("a:gd")):
            name = gd.get("name")
            fmla = gd.get("fmla", "")

            if fmla.startswith("val "):
                value = int(fmla.split()[1])
                result[name] = value

    except Exception:
        pass

    return result


def calc_callout_pointer(shape, geom: str, adjustments: dict) -> dict | None:
    """
    计算吹き出し(callout)尾巴指向点坐标
    返回:
    {
        "x": xxx,
        "y": xxx
    }
    """

    try:
        left = emu_to_pt(shape.left)
        top = emu_to_pt(shape.top)
        width = emu_to_pt(shape.width)
        height = emu_to_pt(shape.height)

        if None in (left, top, width, height):
            return None

        # Office adjustment 参数
        adj1 = adjustments.get("adj1")
        adj2 = adjustments.get("adj2")

        if adj1 is None or adj2 is None:
            return None

        x = left + width * (adj1 / 100000)
        y = top + height * (adj2 / 100000)

        return {
            "x": round(x, 2),
            "y": round(y, 2)
        }

    except Exception:
        return None


def parse_shape(shape) -> dict:
    """解析普通图形（矩形、圆形、自定义等）"""
    info: dict[str, Any] = {
        "id":   shape.shape_id,
        "name": shape.name,
        "type": "shape",
        "bbox": bbox(shape),
    }
    geom = get_shape_geometry(shape)
    if geom:
        info["geometry"] = geom

        adjustments = get_shape_adjustments(shape)
        if adjustments and "callout" in geom.lower():

            pointer = calc_callout_pointer(
                shape,
                geom,
                adjustments
            )

            if pointer:
                info["pointer"] = pointer
    try:
        if shape.has_text_frame:
            paras = parse_text_frame(shape.text_frame)
            if paras:
                info["text"] = paras
                info["text_plain"] = text_plain(shape.text_frame)
    except Exception:
        pass
    return info


# ─────────────────────────── 分组递归 ───────────────────────────
def parse_group(group_shape, shape_id_map: dict, z_index: int = 0) -> dict:
    """递归解析分组，返回组内元素列表"""
    info: dict[str, Any] = {
        "id":   group_shape.shape_id,
        "name": group_shape.name,
        "type": "group",
        "bbox": bbox(group_shape),
        "children": [],
    }
    for i, child in enumerate(group_shape.shapes):
        parsed = dispatch_shape(child, shape_id_map, z_index=i)
        if parsed:
            info["children"].append(parsed)
    return info


# ─────────────────────────── 形状分发 ───────────────────────────
def dispatch_shape(shape, shape_id_map: dict, z_index: int = 0) -> dict | None:
    """根据形状类型分发到对应解析函数"""
    try:
        st = shape.shape_type
    except Exception:
        st = None

    result: dict | None = None

    if is_connector(shape):
        result = parse_connector(shape, shape_id_map)
    elif st == MSO_SHAPE_TYPE.TABLE:
        result = parse_table(shape)
    elif st == MSO_SHAPE_TYPE.PICTURE:
        result = parse_picture(shape)
    elif st == MSO_SHAPE_TYPE.GROUP:
        result = parse_group(shape, shape_id_map, z_index)
    else:
        result = parse_shape(shape)

    if result is not None:
        result["z_index"] = z_index

    return result


# ─────────────────────────── 层叠关系分析 ───────────────────────────
def analyze_overlap(elements: list[dict]) -> list[dict]:
    """
    分析非连接线元素之间的覆盖关系。
    为每个元素添加 covered_by / covers 字段（shape_id 列表）。
    """
    # 只考虑有实体位置的元素（非连接线）
    solid = [e for e in elements if e.get("type") not in ("connector",) and e.get("bbox")]

    for i, elem in enumerate(solid):
        covered_by = []
        covers = []
        bb_i = elem["bbox"]
        zi = elem.get("z_index", i)

        for j, other in enumerate(solid):
            if other["id"] == elem["id"]:
                continue
            bb_j = other["bbox"]
            zj = other.get("z_index", j)
            if bbox_overlaps(bb_i, bb_j):
                if zj > zi:
                    covered_by.append(other["id"])
                else:
                    covers.append(other["id"])

        if covered_by:
            elem["covered_by"] = covered_by  # 被这些 shape_id 覆盖
        if covers:
            elem["covers"] = covers          # 覆盖了这些 shape_id

    return elements


# ─────────────────────────── 连接线逻辑图 ───────────────────────────
def build_connection_graph(elements: list[dict]) -> list[dict]:
    """
    从连接线中提取逻辑关系，生成 edges 列表。
    edge: { from_id, from_name, to_id, to_name, label?, arrows? }
    """
    id_map = {e["id"]: e for e in elements}
    edges = []

    for elem in elements:
        if elem.get("type") != "connector":
            continue

        edge: dict[str, Any] = {
            "connector_id": elem["id"],
            "connector_name": elem["name"],
        }
        if "from_shape" in elem:
            fs = elem["from_shape"]
            edge["from_id"]   = fs["shape_id"]
            edge["from_name"] = fs.get("shape_name") or id_map.get(fs["shape_id"], {}).get("name")
        if "to_shape" in elem:
            ts = elem["to_shape"]
            edge["to_id"]   = ts["shape_id"]
            edge["to_name"] = ts.get("shape_name") or id_map.get(ts["shape_id"], {}).get("name")
        if "arrows" in elem:
            edge["arrows"] = elem["arrows"]
        if "label" in elem:
            edge["label"] = elem["label"]

        edges.append(edge)

    return edges


# ─────────────────────────── 幻灯片解析 ───────────────────────────
def parse_slide(slide, slide_index: int) -> dict:
    """解析单张幻灯片"""
    # 第一遍：建立 shape_id → {name} 映射（用于连接线引用）
    shape_id_map: dict[int, dict] = {}
    for shape in slide.shapes:
        shape_id_map[shape.shape_id] = {"name": shape.name}

    # 第二遍：逐形状解析（保留 z-order）
    elements: list[dict] = []
    for z_idx, shape in enumerate(slide.shapes):
        parsed = dispatch_shape(shape, shape_id_map, z_index=z_idx)
        if parsed:
            elements.append(parsed)

    # 层叠关系分析
    analyze_overlap(elements)

    # 连接线逻辑图
    edges = build_connection_graph(elements)

    # 幻灯片标题（第一个占位符文本）
    title = None
    try:
        title = slide.shapes.title.text if slide.shapes.title else None
    except Exception:
        pass

    return {
        "slide_index": slide_index,
        "title": title,
        "elements": elements,
        "connection_graph": {
            "description": "图形间逻辑指向关系（edges）",
            "edges": edges,
        },
    }


# ─────────────────────────── 主入口 ───────────────────────────
def parse_pptx(file_path: str) -> dict:
    """解析整个 PPTX 文件，返回结构化 JSON dict"""
    prs = Presentation(file_path)
    result: dict[str, Any] = {
        "file": Path(file_path).name,
        "slide_width_pt":  emu_to_pt(prs.slide_width),
        "slide_height_pt": emu_to_pt(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": [],
    }
    for i, slide in enumerate(prs.slides):
        result["slides"].append(parse_slide(slide, slide_index=i + 1))
    return result


def to_markdown_summary(data: dict) -> str:
    """
    将解析结果转换为可读 Markdown 摘要（附带完整 JSON 代码块）。
    """
    lines = [
        f"# PPTX 解析结果：{data['file']}",
        f"",
        f"- 幻灯片尺寸：{data['slide_width_pt']} × {data['slide_height_pt']} pt",
        f"- 幻灯片数量：{data['slide_count']}",
        f"",
    ]
    for slide in data["slides"]:
        lines.append(f"## 第 {slide['slide_index']} 页{('：' + slide['title']) if slide['title'] else ''}")
        lines.append("")

        # 元素汇总
        types: dict[str, int] = {}
        for e in slide["elements"]:
            t = e.get("type", "unknown")
            types[t] = types.get(t, 0) + 1
        lines.append("**元素统计：** " + "、".join(f"{t}×{n}" for t, n in types.items()))
        lines.append("")

        # 连接关系
        edges = slide["connection_graph"]["edges"]
        if edges:
            lines.append("**图形逻辑连接关系：**")
            lines.append("")
            for e in edges:
                frm  = e.get("from_name") or f"id={e.get('from_id','?')}"
                to   = e.get("to_name")   or f"id={e.get('to_id','?')}"
                lbl  = f" [{e['label']}]" if e.get("label") else ""
                arr  = e.get("arrows", {})
                arrow_str = ""
                if arr.get("head") and arr.get("head") != "none":
                    arrow_str = " →"
                elif arr.get("tail") and arr.get("tail") != "none":
                    arrow_str = " ←"
                lines.append(f"- `{frm}`{arrow_str} `{to}`{lbl}")
            lines.append("")

        # 覆盖关系
        covered = [e for e in slide["elements"] if e.get("covered_by")]
        if covered:
            lines.append("**图形覆盖关系（被覆盖元素）：**")
            lines.append("")
            id_name = {e["id"]: e.get("name", str(e["id"])) for e in slide["elements"]}
            for e in covered:
                cover_names = [id_name.get(cid, str(cid)) for cid in e["covered_by"]]
                lines.append(f"- `{e['name']}` (type={e['type']}) 被 {', '.join('`'+n+'`' for n in cover_names)} 覆盖")
            lines.append("")

    # 完整 JSON 
    lines.append("---")
    lines.append("")
    lines.append("## 完整结构化 JSON")
    lines.append("")
    lines.append("```json")
    lines.append(json.dumps(data, ensure_ascii=False, indent=2))
    lines.append("```")

    return "\n".join(lines)


# ─────────────────────────── CLI ───────────────────────────

def get_pptx_source(output_md: str = "./pptx_summary.md") -> dict:
    pptx_path = "C:/Users/tR16277/Desktop/work/work06/file/a.pptx"
    data = parse_pptx(pptx_path)
    r = to_markdown_summary(data)
    with open(output_md, "w", encoding="utf-8") as f:
        f.write(r)
    return data
