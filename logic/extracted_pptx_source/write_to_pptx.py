from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def rebuild_shape(slide, shape_data):

    left = shape_data["left"]
    top = shape_data["top"]
    width = shape_data["width"]
    height = shape_data["height"]

    text = shape_data.get("text", "")

    shape_type = shape_data["type"]

    # 文本框
    if "TEXT_BOX" in shape_type:

        textbox = slide.shapes.add_textbox(
            left,
            top,
            width,
            height
        )

        textbox.text = text

    # rectangle
    elif "RECTANGLE" in shape_type:

        rect = slide.shapes.add_textbox(
            left,
            top,
            width,
            height
        )

        rect.text = text

    # line
    elif "LINE" in shape_type:

        slide.shapes.add_connector(
            1,
            left,
            top,
            left + width,
            top + height
        )


def rebuild_ppt(data):

    prs = Presentation()

    blank_layout = prs.slide_layouts[6]

    for slide_data in data:

        slide = prs.slides.add_slide(blank_layout)

        for shape_data in slide_data["shapes"]:

            rebuild_shape(slide, shape_data)

    prs.save("output/rebuilt.pptx")


if __name__ == "__main__":

    import json

    with open(
        "./slide_data.json",
        encoding="utf-8"
    ) as f:

        data = json.load(f)

    rebuild_ppt(data)

    print("生成完成")