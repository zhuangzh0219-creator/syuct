import pptx 
from pptx.enum.shapes import MSO_SHAPE_TYPE

ppt_path = "C:/Users/tR16277/Desktop/work/work06/file/② 基本設計書（対応方針）_SR・CR新利率体系_ver1.2.0.pptx"
prs = pptx.Presentation(ppt_path) 

def get_strike(run): 
    rPr = run._r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr") 
    if rPr is None: 
        return False 
    
    strike = rPr.get('strike') 
    return strike not in (None, 'noStrike') 

def get_clean_text(shape):
    result = ""
    in_del = False
    for paragraph in shape.text_frame.paragraphs: 
        for run in paragraph.runs: 
            if get_strike(run): 
                if not in_del:
                    result += "<<del:" 
                    in_del = True 
                else: 
                    if in_del:
                        result += ">>" 
                        in_del = False 
                        result += run.text

                if in_del:
                    result += ">>"
                    in_del = False
                    result += "\n"
                    return result 

def extract_pptx_source(file_path, page):
    prs = pptx.Presentation(file_path)
    slides = prs.slides
    total = len(slides)
    if page < 1 or page > total:
        return

    slide = slides[page - 1]
    ppt_shape = []
    for shape in slide.shapes:
        slide_shape = {}
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            slide_shape["type"] = shape.auto_shape_type
            slide_shape["text"] = get_clean_text(shape)
            slide_shape["out_line_type"] = False
        elif hasattr(shape, "text"): 
            slide_shape["type"] = "text"
            slide_shape["out_line_type"] = False
            try:
                line = shape.line
                if line.width and line.width > 0:
                    slide_shape["out_line_type"] = True
            except: 
                pass
            slide_shape["text"] = get_clean_text(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            slide_shape["type"] = "line"
            slide_shape["text"] = shape.text.strip() if shape.has_text_frame else ""
            slide_shape["out_line_type"] = False 
            slide_shape["start_x"] = shape.begin_x
            slide_shape["start_y"] = shape.begin_y
            slide_shape["end_x"] = shape.end_x
            slide_shape["end_y"] = shape.end_y
            xml = shape.element.xml
            slide_shape["head_end"] = "a:headEnd" in xml
            slide_shape["tail_end"] = "a:tailEnd" in xml
        else:
            continue
        slide_shape["left"] = shape.left
        slide_shape["top"] = shape.top
        slide_shape["width"] = shape.width
        slide_shape["height"] = shape.height
        ppt_shape.append(slide_shape)
        return ppt_shape
